import csv
import pathlib
import random
from dataclasses import dataclass
from typing import List

import pptx
import typer
import pylightxl as xl


@dataclass
class Question:
    question: str
    answer: str


def create_presentation(input_template: pathlib.Path,
                        slide_name: str,
                        question_input: pathlib.Path,
                        output: pathlib.Path,
                        seed: int = 42):
    random.seed(seed)

    def match_questions(questions_a: List[Question]):
        # Make a copy of the list to avoid modifying the original
        questions_b = list(questions_a)

        # Iterate over the shuffled list in pairs
        invalid = True
        matched_a, matched_b = [], []

        while invalid:
            # Shuffle the list randomly
            random.shuffle(questions_b)

            for i in range(len(questions_a)):
                # Check if the pair contains two identical questions_a
                if questions_a[i].question == questions_b[i].question:
                    # If it does, swap one of them with the next element in the list until a valid pair is found
                    for j in range(i + 1, len(questions_b)):
                        if questions_a[i].question != questions_b[j].question:
                            questions_b[i], questions_b[j] = questions_b[j], questions_b[i]
                            break

                # Add the valid pair to the list of pairs
                matched_a += [questions_a[i]]
                matched_b += [questions_b[i]]

            # small validity check
            assert len(matched_a) == len(questions_a)
            for i in range(len(matched_a)):
                if matched_a[0] != matched_b[1]:
                    invalid = False
                else:
                    invalid = True
                    matched_a, matched_b = [], []
                    break

        return matched_a, matched_b

    input_data = []

    if not str(input_template).endswith(".pptx"):
        print("The input template is not a .pptx file...")
        raise typer.Exit(code=1)
    if not str(output).endswith(".pptx"):
        print("The output path is not a .pptx file...")
        raise typer.Exit(code=1)
    if not str(question_input).endswith(".xlsx"):
        print("The question template is not a .xlsx file...")
        raise typer.Exit(code=1)

    if not input_template.exists():
        print("The input template does not exist...")
        raise typer.Exit(code=1)

    if not question_input.exists():
        print("The question template does not exist...")
        raise typer.Exit(code=1)


    db = xl.readxl(question_input)
    selected_table = db.ws_names[0]
    for row in db.ws(ws=selected_table).rows:
        input_data += [Question(row[0], row[1])]

    pairs_a, pairs_b = match_questions(input_data)
    prs = pptx.Presentation(input_template)
    slide_layout = prs.slide_layouts.get_by_name(slide_name)
    for i in range(len(pairs_a)):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Frage {i + 1}"
        list(slide.placeholders)[1].text = pairs_a[i].question
        list(slide.placeholders)[2].text = pairs_b[i].question
    prs.save(output)

    with open(str(output) + ".txt", 'w+') as file:
        file.writelines("Lösungen Gruppe A\n")
        file.writelines([f"Frage {i + 1}: {question.answer}\n" for i, question in enumerate(pairs_a)])

        file.writelines(["\n", "\n", "Lösungen Gruppe B\n"])
        file.writelines([f"Frage {i + 1}: {question.answer}\n" for i, question in enumerate(pairs_b)])


if __name__ == '__main__':
    typer.run(create_presentation)
